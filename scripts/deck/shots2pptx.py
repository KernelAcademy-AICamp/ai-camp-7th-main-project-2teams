"""촬영한 슬라이드 이미지 → PPTX.

텍스트를 PPT 도형으로 재조립하는 방식은 HTML 디자인과 어긋난다.
렌더 결과를 한 장씩 그대로 깔면 폰트·색·여백이 원본과 일치한다.
대신 PPT 안에서 텍스트 편집은 불가능하다 — 수정은 HTML에서 하고 다시 뽑는다.

사용: python3 scripts/deck/shots2pptx.py   (scripts/deck/build.sh가 호출)
필요: pip install python-pptx
"""

import html
import re
from pathlib import Path

from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from pptx.oxml.ns import qn
from pptx.util import Inches

ROOT = Path(__file__).resolve().parents[2]
DOCS = ROOT / "docs"
SHOTS = DOCS / ".deck-shots"

SLIDE_W = Inches(13.333)  # 16:9 — 촬영 뷰포트 1600x900과 같은 비율이라 왜곡이 없다
SLIDE_H = Inches(7.5)

# (촬영 디렉터리, 산출 PPTX, 대본을 가져올 HTML)
DECKS = [
    ("sub", "service-plan-presentation.pptx", None),
    ("talk", "service-plan-presentation-talk.pptx", "service-plan-presentation-talk.html"),
]


def extract_notes(deck: Path) -> list[str]:
    """발표용 덱의 .note 대본을 슬라이드 순서대로 뽑는다. 없는 슬라이드는 빈 문자열."""
    chunks = re.split(r'<section class="slide', deck.read_text())[1:]
    notes = []
    for chunk in chunks:
        found = re.search(r'<p class="note">(.*?)</p>', chunk, re.S)
        text = re.sub(r"<[^>]+>", "", found.group(1)) if found else ""
        notes.append(" ".join(html.unescape(text).split()))
    return notes



def declare_notes_master(prs) -> None:
    """presentation.xml에 notesMasterIdLst를 넣는다.

    노트를 추가하면 notesMaster 파트는 생기는데 python-pptx는 presentation.xml의
    선언을 채우지 않는다. 파트는 있는데 선언이 없는 상태라 엄격한 파서(Keynote)가
    "파일 포맷이 유효하지 않다"로 거부한다.
    """
    pres_elm = prs._element
    if pres_elm.find(qn("p:notesMasterIdLst")) is not None:
        return
    notes_master_part = prs.notes_master.part
    r_id = prs.part.relate_to(notes_master_part, RT.NOTES_MASTER)
    id_lst = pres_elm.makeelement(qn("p:notesMasterIdLst"), {})
    id_lst.append(id_lst.makeelement(qn("p:notesMasterId"), {qn("r:id"): r_id}))
    # 스키마 순서: sldMasterIdLst → notesMasterIdLst → ... 순서를 어기면 그것도 거부 사유다.
    pres_elm.find(qn("p:sldMasterIdLst")).addnext(id_lst)


def build(shots_dir: Path, out: Path, notes: list[str]) -> None:
    images = sorted(shots_dir.glob("slide-*.jpg"))
    if not images:
        raise SystemExit(f"촬영 이미지 없음: {shots_dir} — build.sh로 촬영부터 실행")

    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    # 기본 템플릿이 남긴 type="screen4x3"를 지운다. 크기는 16:9인데 선언만 4:3이라
    # 엄격한 파서(Keynote 등)가 "파일 포맷이 유효하지 않다"로 거부한다.
    sld_sz = prs._element.find(qn("p:sldSz"))
    if sld_sz is not None and sld_sz.get("type"):
        del sld_sz.attrib["type"]
    if any(notes):
        declare_notes_master(prs)
    blank = prs.slide_layouts[6]

    for i, image in enumerate(images):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(image), 0, 0, width=SLIDE_W, height=SLIDE_H)
        if i < len(notes) and notes[i]:
            slide.notes_slide.notes_text_frame.text = notes[i]

    prs.save(out)
    filled = sum(1 for n in notes[: len(images)] if n)
    print(f"{out.name}: {len(images)}장 · 발표자 노트 {filled}개 · {out.stat().st_size // 1024}KB")


if __name__ == "__main__":
    for shots_name, out_name, notes_html in DECKS:
        build(
            SHOTS / shots_name,
            DOCS / out_name,
            extract_notes(DOCS / notes_html) if notes_html else [],
        )
